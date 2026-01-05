"""
Export Orchestrator - Main coordinator for PowerPoint export functionality.
Generates presentation slides with the same content as the PDF export.
Supports both programmatic generation and screenshot-based capture.
"""

import os
import json
import logging
import re
from io import BytesIO
from typing import Optional, Dict, Any, List
from datetime import datetime

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Configure logging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(name)s: %(message)s')

# Screenshot service imports (optional - gracefully handle if not available)
try:
    from core.streamlit_screenshot import capture_streamlit_screenshots
    SCREENSHOT_AVAILABLE = True
except ImportError:
    SCREENSHOT_AVAILABLE = False
    logger.warning("Screenshot service not available")

# Metric descriptions for fallback
METRICS = {
    "Representation": {
        "high": "Strong representation of diverse perspectives and communities.",
        "medium": "Adequate representation with room for improvement.",
        "low": "Limited representation; consider expanding diverse voices."
    },
    "Cultural Relevance": {
        "high": "Excellent cultural alignment with target audiences.",
        "medium": "Good cultural connections with some gaps.",
        "low": "Limited cultural relevance; needs deeper cultural insight."
    },
    "Platform Relevance": {
        "high": "Strong platform strategy aligned with audience behavior.",
        "medium": "Moderate platform alignment with optimization opportunities.",
        "low": "Platform strategy needs significant refinement."
    },
    "Cultural Vernacular": {
        "high": "Authentic use of language that resonates with audiences.",
        "medium": "Generally appropriate language with room for refinement.",
        "low": "Language may not fully connect with target audiences."
    },
    "Media Ownership Equity": {
        "high": "Strong investment in diverse media channels.",
        "medium": "Moderate diversity in media investment.",
        "low": "Consider expanding to more diverse-owned media."
    },
    "Cultural Authority": {
        "high": "Strong credibility and authority within cultural context.",
        "medium": "Building authority with room for growth.",
        "low": "Limited cultural authority; needs authentic partnerships."
    },
    "Buzz & Conversation": {
        "high": "High potential for organic conversation and sharing.",
        "medium": "Moderate buzz potential with optimization needed.",
        "low": "Limited conversation potential; needs creative refresh."
    },
    "Commerce Bridge": {
        "high": "Strong connection between awareness and purchase intent.",
        "medium": "Moderate commerce bridge with conversion opportunities.",
        "low": "Weak link between messaging and purchase action."
    },
    "Geo-Cultural Fit": {
        "high": "Excellent geographic and cultural market alignment.",
        "medium": "Good regional fit with some expansion opportunities.",
        "low": "Geographic targeting needs refinement."
    }
}


class ExportOrchestrator:
    """
    Orchestrates the complete export process generating native PowerPoint
    slides with the same data and structure as the PDF export.
    """

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Colors (matching PDF)
    PRIMARY_COLOR = RGBColor(59, 130, 246)    # #3b82f6 - blue for QVI
    SECONDARY_COLOR = RGBColor(5, 150, 105)   # #059669 - green for streaming
    TEXT_COLOR = RGBColor(0, 0, 0)            # black
    GRAY = RGBColor(100, 116, 139)            # #64748b
    LIGHT_BLUE = RGBColor(219, 234, 254)      # #dbeafe
    LIGHT_GREEN = RGBColor(209, 250, 229)     # #d1fae5
    WHITE = RGBColor(255, 255, 255)
    PURPLE = RGBColor(97, 113, 234)           # #6171EA - cover bg

    # Font
    FONT_NAME = "Helvetica"

    def __init__(self, session_state: Dict[str, Any]):
        """Initialize orchestrator with Streamlit session state."""
        self.session_state = session_state

    def export_presentation(
        self,
        brand_name: str = "Brand",
        industry: str = "General",
        progress_callback: Optional[callable] = None
    ) -> bytes:
        """Generate complete PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        def update_progress(percent: int, message: str):
            if progress_callback:
                progress_callback(min(percent, 100), message)

        try:
            logger.info("=" * 70)
            logger.info("POWERPOINT EXPORT STARTED")
            logger.info(f"Brand: {brand_name} | Industry: {industry}")
            logger.info("=" * 70)

            # 1. Cover slide
            update_progress(10, "Creating cover slide...")
            logger.info("Creating cover slide...")
            self._add_cover_slide(prs, brand_name, industry)

            # 2. Metric Breakdown slide
            update_progress(15, "Creating metrics slide...")
            logger.info("Creating metric breakdown slide...")
            self._add_metric_breakdown_slide(prs)

            # 3. Media Affinities (all in one slide)
            update_progress(30, "Creating media affinities slide...")
            logger.info("Creating media affinities slide...")
            self._add_media_affinities_combined_slide(prs)

            # 4. Audience Insights slide (with psychographic highlights)
            update_progress(45, "Creating audience insights slide...")
            logger.info("Creating audience insights slide...")
            self._add_audience_insights_slide(prs)

            # 5. Competitor Tactics slide (only if data exists)
            competitor_tactics = self.session_state.get('competitor_tactics', [])
            if competitor_tactics and isinstance(competitor_tactics, list) and len(competitor_tactics) > 0:
                update_progress(55, "Creating competitor tactics slide...")
                logger.info("Creating competitor tactics slide...")
                self._add_competitor_tactics_slide(prs)
            else:
                logger.info("Skipping competitor tactics slide (no data)")

            # 6. Marketing Trends slide (includes benchmark at bottom)
            update_progress(80, "Creating marketing trends slide...")
            logger.info("Creating marketing trends slide...")
            self._add_marketing_trends_slide(prs)

            # # 7. Footer slide
            # update_progress(90, "Finalizing presentation...")
            # logger.info("Adding footer slide...")
            # self._add_footer_slide(prs)

            # Save to bytes
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            file_size = len(output.getvalue())
            slide_count = len(prs.slides)
            logger.info("=" * 70)
            logger.info("POWERPOINT EXPORT COMPLETED")
            logger.info(f"Total slides: {slide_count}")
            logger.info(f"File size: {file_size:,} bytes ({file_size/1024:.1f} KB)")
            logger.info("=" * 70)

            return output.getvalue()

        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise

    def _add_cover_slide(self, prs: Presentation, brand_name: str, industry: str):
        """Add cover slide with purple background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Purple background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_WIDTH, self.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.PURPLE
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Audience Resonance Index Scorecard"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{brand_name} | {industry}"
        p.font.size = Pt(20)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%B %d, %Y")
        p.font.size = Pt(14)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

    def _add_metric_breakdown_slide(self, prs: Presentation):
        """Add detailed metrics table slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, "Metric Breakdown")

        scores = self.session_state.get('scores', {})
        ai_insights = self.session_state.get('ai_insights', {})
        metric_details = ai_insights.get('metric_details', {}) if ai_insights else {}

        if not scores:
            return

        # Table header
        y_pos = Inches(1.2)
        header_height = Inches(0.4)

        # Header background
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.3), y_pos,
            Inches(12.7), header_height
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.LIGHT_BLUE
        header_bg.line.fill.background()

        # Header text
        headers = [("Metric", 0.4, 2.0), ("Score", 2.5, 0.8), ("Description", 3.4, 9.5)]
        for text, left, width in headers:
            box = slide.shapes.add_textbox(Inches(left), y_pos + Inches(0.08), Inches(width), header_height)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

        y_pos += header_height
        row_height = Inches(0.55)

        # Data rows
        for i, (metric, score) in enumerate(scores.items()):
            # Alternating row colors
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.3), y_pos,
                Inches(12.7), row_height
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(224, 237, 255) if i % 2 == 0 else RGBColor(240, 249, 255)
            row_bg.line.fill.background()

            # Metric name
            box = slide.shapes.add_textbox(Inches(0.4), y_pos + Inches(0.12), Inches(2.0), row_height)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = metric
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            # Score
            box = slide.shapes.add_textbox(Inches(2.5), y_pos + Inches(0.12), Inches(0.8), row_height)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{score}/10"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            # Description
            level = "high" if score >= 7 else "medium" if score >= 4 else "low"
            description = metric_details.get(metric, METRICS.get(metric, {}).get(level, ""))

            box = slide.shapes.add_textbox(Inches(3.4), y_pos + Inches(0.05), Inches(9.5), row_height)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description[:180] if len(description) > 180 else description
            p.font.size = Pt(8)
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.GRAY

            y_pos += row_height

    def _add_media_affinities_combined_slide(self, prs: Presentation):
        """Add media affinities, TV networks, and streaming on ONE slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, "Media Affinities")

        # Get all data
        media_affinity = self.session_state.get('media_affinity', '[]')
        if isinstance(media_affinity, str):
            try:
                media_affinity = json.loads(media_affinity)
            except:
                media_affinity = []

        audience_media = self.session_state.get('audience_media_consumption', {})
        if isinstance(audience_media, str):
            try:
                audience_media = json.loads(audience_media)
            except:
                audience_media = {}

        tv_networks = audience_media.get('tv_networks', [])
        streaming_platforms = audience_media.get('streaming_platforms', [])

        y_pos = Inches(1.1)

        # === TOP MEDIA AFFINITY SITES ===
        sub_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12), Inches(0.35))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Top Media Affinity Sites"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.35)

        # QVI explanation
        qvi_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12), Inches(0.25))
        tf = qvi_box.text_frame
        p = tf.paragraphs[0]
        p.text = "QVI = Quality Visit Index, a score indicating audience engagement strength"
        p.font.size = Pt(8)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.GRAY

        y_pos += Inches(0.3)

        # Default card height for layout calculations
        card_height = Inches(0.7)

        # Media sites - compact grid (5 per row)
        if media_affinity:
            x_positions = [Inches(0.3), Inches(2.85), Inches(5.4), Inches(7.95), Inches(10.5)]
            card_width = Inches(2.4)

            for i, site in enumerate(media_affinity[:5]):
                x = x_positions[i]

                # Card background
                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, card_width, card_height)
                card.fill.solid()
                card.fill.fore_color.rgb = self.LIGHT_BLUE
                card.line.fill.background()

                # Text
                text_box = slide.shapes.add_textbox(x + Inches(0.08), y_pos + Inches(0.08), card_width - Inches(0.16), card_height)
                tf = text_box.text_frame
                tf.word_wrap = True

                p1 = tf.paragraphs[0]
                p1.text = site.get('name', 'Unknown')
                p1.font.size = Pt(9)
                p1.font.bold = True
                p1.font.name = self.FONT_NAME
                p1.font.color.rgb = self.TEXT_COLOR

                p2 = tf.add_paragraph()
                p2.text = site.get('category', '')
                p2.font.size = Pt(7)
                p2.font.name = self.FONT_NAME
                p2.font.color.rgb = self.GRAY

                p3 = tf.add_paragraph()
                p3.text = f"QVI: {site.get('qvi', 0)}"
                p3.font.size = Pt(8)
                p3.font.bold = True
                p3.font.name = self.FONT_NAME
                p3.font.color.rgb = self.PRIMARY_COLOR

        y_pos += card_height + Inches(0.25)

        # === TOP TV NETWORK AFFINITIES ===
        if tv_networks:
            sub_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12), Inches(0.35))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Top TV Network Affinities"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.4)

            x_positions = [Inches(0.3), Inches(2.85), Inches(5.4), Inches(7.95), Inches(10.5)]
            card_width = Inches(2.4)
            card_height = Inches(0.7)

            for i, network in enumerate(tv_networks[:5]):
                x = x_positions[i]

                card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, card_width, card_height)
                card.fill.solid()
                card.fill.fore_color.rgb = self.LIGHT_BLUE
                card.line.fill.background()

                text_box = slide.shapes.add_textbox(x + Inches(0.08), y_pos + Inches(0.08), card_width - Inches(0.16), card_height)
                tf = text_box.text_frame

                p1 = tf.paragraphs[0]
                p1.text = network.get('name', '')
                p1.font.size = Pt(9)
                p1.font.bold = True
                p1.font.name = self.FONT_NAME

                p2 = tf.add_paragraph()
                p2.text = network.get('category', '')
                p2.font.size = Pt(7)
                p2.font.name = self.FONT_NAME
                p2.font.color.rgb = self.GRAY

                p3 = tf.add_paragraph()
                p3.text = f"QVI: {network.get('qvi', 0)}"
                p3.font.size = Pt(8)
                p3.font.bold = True
                p3.font.name = self.FONT_NAME
                p3.font.color.rgb = RGBColor(30, 136, 229)  # #1e88e5

            y_pos += card_height + Inches(0.25)

        # === TOP STREAMING PLATFORMS ===
        if streaming_platforms:
            sub_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12), Inches(0.35))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Top Streaming Platforms"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.4)

            # 3 per row for streaming
            x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
            card_width = Inches(3.9)
            card_height = Inches(0.65)

            for row in range(2):
                for col in range(3):
                    idx = row * 3 + col
                    if idx >= len(streaming_platforms):
                        break

                    platform = streaming_platforms[idx]
                    x = x_positions[col]
                    y = y_pos + (row * (card_height + Inches(0.1)))

                    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
                    card.fill.solid()
                    card.fill.fore_color.rgb = self.LIGHT_GREEN
                    card.line.fill.background()

                    text_box = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.08), card_width - Inches(0.16), card_height)
                    tf = text_box.text_frame

                    p1 = tf.paragraphs[0]
                    p1.text = platform.get('name', '')
                    p1.font.size = Pt(9)
                    p1.font.bold = True
                    p1.font.name = self.FONT_NAME

                    p2 = tf.add_paragraph()
                    p2.text = f"{platform.get('category', '')}   QVI: {platform.get('qvi', 0)}"
                    p2.font.size = Pt(8)
                    p2.font.name = self.FONT_NAME
                    p2.font.color.rgb = self.SECONDARY_COLOR

    def _add_audience_insights_slide(self, prs: Presentation):
        """Add audience insights slide with psychographic highlights and segments."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, "Audience Insights")

        y_pos = Inches(0.75)

        # Psychographic Highlights section first
        psycho_text = self.session_state.get('pychographic_highlights', '')
        if psycho_text:
            # Section header
            header_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
            tf = header_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Psychographic Highlights"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.3)

            # Clean HTML tags
            clean_text = re.sub(r'<[^>]+>', '', psycho_text)

            # Psychographic content
            content_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(1.2))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = clean_text
            p.font.size = Pt(9)
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.5)

        # Audience Summary section
        header_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Audience Summary"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.35)

        audience_segments = self.session_state.get('audience_segments', {})
        segments = audience_segments.get('segments', []) if isinstance(audience_segments, dict) else []

        # Audience types
        audience_labels = [
            "Primary Growth Audience:",
            "Secondary Growth Audience:",
            "Emerging Audience Opportunity:"
        ]

        for i, label in enumerate(audience_labels):
            if i >= len(segments):
                break

            segment = segments[i]

            # Section header (bold)
            header_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
            tf = header_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.3)

            # Segment info
            name = segment.get('name', '')
            description = segment.get('description', '')

            info_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.35))
            tf = info_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{name}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.25)

            # Demographics
            targeting = segment.get('targeting_params', {})
            if targeting:
                demo_parts = []
                if targeting.get('age_range'):
                    demo_parts.append(f"Age: {targeting['age_range']}")
                if targeting.get('gender_targeting'):
                    demo_parts.append(f"Gender: {targeting['gender_targeting']}")
                if targeting.get('income_targeting'):
                    demo_parts.append(f"Income: {targeting['income_targeting']}")

                if demo_parts:
                    demo_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.25))
                    tf = demo_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = "Demographics: " + " | ".join(demo_parts)
                    p.font.size = Pt(9)
                    p.font.name = self.FONT_NAME
                    p.font.color.rgb = self.TEXT_COLOR
                    y_pos += Inches(0.25)

            # Platform recommendation
            platform_targeting = segment.get('platform_targeting', [])
            if platform_targeting:
                platform = platform_targeting[0].get('platform', '')
                if platform:
                    # Determine metric type based on platform
                    performance = segment.get('expected_performance', {})
                    metric_value = performance.get('CTR', '0.05-0.7%')
                    metric_name = "Expected CTR"

                    if 'video' in platform.lower() or 'ott' in platform.lower() or 'ctv' in platform.lower():
                        metric_name = "Expected VCR"
                        metric_value = "90-100%"
                    elif 'audio' in platform.lower() or 'podcast' in platform.lower():
                        metric_name = "Expected LTR"
                        metric_value = "80-95%"

                    plat_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.25))
                    tf = plat_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"Recommended Platform: {platform} ({metric_name}: {metric_value})"
                    p.font.size = Pt(9)
                    p.font.name = self.FONT_NAME
                    p.font.color.rgb = self.TEXT_COLOR
                    y_pos += Inches(0.25)

            # Interests (for emerging audience)
            if i == 2:  # Emerging audience
                interests = segment.get('interest_categories', [])
                if interests:
                    int_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.25))
                    tf = int_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = "Key Interests: " + ", ".join(interests)
                    p.font.size = Pt(9)
                    p.font.name = self.FONT_NAME
                    p.font.color.rgb = self.TEXT_COLOR
                    y_pos += Inches(0.25)

                rationale = segment.get('rationale', '')
                if rationale:
                    rat_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.5))
                    tf = rat_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"Rationale: {rationale}"
                    p.font.size = Pt(9)
                    p.font.name = self.FONT_NAME
                    p.font.color.rgb = self.TEXT_COLOR
                    y_pos += Inches(0.5)

            y_pos += Inches(0.2)

    def _add_competitor_tactics_slide(self, prs: Presentation):
        """Add competitor tactics slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, "Competitor Tactics")

        competitor_tactics = self.session_state.get('competitor_tactics', [])

        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(12.7), Inches(6))
        tf = content_box.text_frame
        tf.word_wrap = True

        if competitor_tactics and isinstance(competitor_tactics, list) and len(competitor_tactics) > 0:
            for i, tactic in enumerate(competitor_tactics[:5]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{i+1}. {tactic}"
                p.font.size = Pt(11)
                p.font.name = self.FONT_NAME
                p.font.color.rgb = self.TEXT_COLOR
                p.space_after = Pt(10)
        else:
            p = tf.paragraphs[0]
            p.text = "Competitor tactics analysis not available. Visit the Competitor Tactics tab to generate custom competitive strategy recommendations."
            p.font.size = Pt(11)
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.GRAY

    def _add_marketing_trends_slide(self, prs: Presentation):
        """Add marketing trends slide - plain text like PDF."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, "Marketing Trend Analysis")

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.75), Inches(12.7), Inches(0.35))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Top Marketing Trends Relevant to Your Campaign"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        # Get trend data
        brief_text = self.session_state.get('brief_text', '')
        try:
            from app.components.marketing_trends import generate_simplified_trend_data
            top_trends, top_markets, top_combinations = generate_simplified_trend_data(brief_text=brief_text)
        except Exception as e:
            logger.warning(f"Could not generate trend data: {e}")
            top_trends = []
            top_markets = []
            top_combinations = []

        y_pos = Inches(1.1)

        # Emerging Trends section
        trends_header = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(6), Inches(0.3))
        tf = trends_header.text_frame
        p = tf.paragraphs[0]
        p.text = "Emerging Trends:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.3)
        for trend in top_trends[:3]:
            if isinstance(trend, dict):
                text = f"• {trend.get('trend', 'Unknown')}: {trend.get('growth', 0)}% growth"
                box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(6), Inches(0.25))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(10)
                p.font.name = self.FONT_NAME
                p.font.color.rgb = self.TEXT_COLOR
                y_pos += Inches(0.25)

        y_pos += Inches(0.15)

        # High-Performance Targets section
        targets_header = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(6), Inches(0.3))
        tf = targets_header.text_frame
        p = tf.paragraphs[0]
        p.text = "High-Performance Targets:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.3)
        for market in top_markets[:3]:
            if isinstance(market, dict):
                text = f"• {market.get('market', 'Unknown')}: {market.get('index', 0)} engagement index"
                box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(6), Inches(0.25))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(10)
                p.font.name = self.FONT_NAME
                p.font.color.rgb = self.TEXT_COLOR
                y_pos += Inches(0.25)

        y_pos += Inches(0.15)

        # Strategic Applications from AI insights
        ai_insights = self.session_state.get('ai_insights', {})
        improvements = ai_insights.get('improvements', []) if ai_insights else []

        if improvements:
            strat_header = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
            tf = strat_header.text_frame
            p = tf.paragraphs[0]
            p.text = "Strategic Applications:"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.3)
            for imp in improvements[:3]:
                if isinstance(imp, dict):
                    area = imp.get('area', '')
                    rec = imp.get('recommendation', '')
                    text = f"• {area}: {rec}"

                    # Calculate dynamic height based on text length
                    # At 9pt Helvetica with 12.7" width, approximately 140 chars per line
                    chars_per_line = 140
                    num_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
                    line_height = 0.15  # inches per line at 9pt
                    box_height = num_lines * line_height + 0.05  # add small padding

                    box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(box_height))
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(9)
                    p.font.name = self.FONT_NAME
                    p.font.color.rgb = self.TEXT_COLOR
                    y_pos += Inches(box_height)

        # Benchmark Comparison section at bottom
        y_pos += Inches(0.2)
        percentile = self.session_state.get('percentile', 50)
        improvement_areas = self.session_state.get('improvement_areas', [])

        # Benchmark header
        bench_header = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
        tf = bench_header.text_frame
        p = tf.paragraphs[0]
        p.text = "Benchmark Comparison"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.3)

        # Benchmark text
        bench_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.8))
        tf = bench_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = (f"This campaign ranks in the top {percentile}% of all campaigns for "
                  f"Audience Resonance Index (ARI), outperforming the majority in relevance, "
                  f"authenticity, and emotional connection.")
        p.font.size = Pt(9)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

        y_pos += Inches(0.25)

        # Improvement areas
        if improvement_areas:
            imp_box = slide.shapes.add_textbox(Inches(0.3), y_pos, Inches(12.7), Inches(0.3))
            tf = imp_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Biggest opportunity areas: {', '.join(improvement_areas[:4])}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = self.FONT_NAME
            p.font.color.rgb = self.TEXT_COLOR

    def _add_footer_slide(self, prs: Presentation):
        """Add footer/closing slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)

        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.3), Inches(1))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "© 2025 Digital Culture Group, LLC. All rights reserved."
        p.font.size = Pt(11)
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.GRAY
        p.alignment = PP_ALIGN.CENTER

    def _set_white_background(self, slide):
        """Set white background on slide."""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_WIDTH, self.SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.WHITE
        bg.line.fill.background()
        # Send to back
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_slide_title(self, slide, title: str):
        """Add consistent title to slide."""
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = self.FONT_NAME
        p.font.color.rgb = self.TEXT_COLOR

    def _add_screenshot_slide(self, prs: Presentation, title: str, image_bytes: bytes):
        """
        Add a slide with a screenshot image.

        Args:
            prs: PowerPoint presentation object
            title: Slide title
            image_bytes: PNG image bytes
        """
        from PIL import Image
        import io

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_white_background(slide)
        self._add_slide_title(slide, title)

        # Load image to get dimensions
        img = Image.open(io.BytesIO(image_bytes))
        img_width, img_height = img.size

        # Calculate scaling to fit slide (leaving room for title)
        # Available area: 13.333" wide x 6.5" tall (after title)
        max_width = Inches(12.7)  # Leave margins
        max_height = Inches(6.3)  # Leave room for title

        # Calculate aspect ratio scaling
        width_ratio = max_width / Emu(img_width * 914400 / 96)  # 96 DPI assumed
        height_ratio = max_height / Emu(img_height * 914400 / 96)
        scale = min(width_ratio, height_ratio, 1.0)  # Don't scale up

        # Calculate final dimensions
        final_width = Emu(img_width * 914400 / 96 * scale)
        final_height = Emu(img_height * 914400 / 96 * scale)

        # Center horizontally
        left = (self.SLIDE_WIDTH - final_width) / 2
        top = Inches(0.9)  # Below title

        # Add image to slide
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(image_stream, left, top, width=final_width, height=final_height)

        logger.info(f"  ✓ Added screenshot slide: {title} ({img_width}x{img_height}px)")

    def export_presentation_with_screenshots(
        self,
        brand_name: str = "Brand",
        industry: str = "General",
        app_url: str = "http://localhost:3006",
        use_live_capture: bool = False,
        progress_callback: Optional[callable] = None
    ) -> bytes:
        """
        Generate PowerPoint presentation using actual screenshots of the Streamlit app.

        Args:
            brand_name: Brand name for the cover slide
            industry: Industry for the cover slide
            app_url: URL of the running Streamlit app
            use_live_capture: If True, capture from running app; if False, render HTML
            progress_callback: Optional callback for progress updates

        Returns:
            PowerPoint file bytes
        """
        if not SCREENSHOT_AVAILABLE:
            logger.warning("Screenshot service not available, falling back to programmatic export")
            return self.export_presentation(brand_name, industry, progress_callback)

        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        def update_progress(percent: int, message: str):
            if progress_callback:
                progress_callback(min(percent, 100), message)

        try:
            logger.info("=" * 70)
            logger.info("POWERPOINT EXPORT WITH SCREENSHOTS STARTED")
            logger.info(f"Brand: {brand_name} | Industry: {industry}")
            logger.info(f"Mode: {'Live capture' if use_live_capture else 'HTML rendering'}")
            logger.info("=" * 70)

            # 1. Cover slide (programmatic - looks better)
            update_progress(5, "Creating cover slide...")
            logger.info("Creating cover slide...")
            self._add_cover_slide(prs, brand_name, industry)

            # 2. Capture Streamlit screenshots
            update_progress(10, "Capturing screenshots...")
            logger.info("Capturing screenshots...")

            screenshots = capture_streamlit_screenshots(
                self.session_state,
                use_live_capture=use_live_capture,
                app_url=app_url
            )

            if not screenshots:
                logger.warning("No screenshots captured, falling back to programmatic export")
                return self.export_presentation(brand_name, industry, progress_callback)

            # 3. Add screenshot slides, inserting Competitor Tactics at position 3
            competitor_tactics = self.session_state.get('competitor_tactics', [])
            total_screenshots = len(screenshots)
            slide_index = 0
            for i, (tab_name, png_bytes) in enumerate(screenshots.items()):
                progress = 20 + int((i / total_screenshots) * 60)
                update_progress(progress, f"Adding slide: {tab_name}...")
                self._add_screenshot_slide(prs, tab_name, png_bytes)
                slide_index += 1

                # Insert Competitor Tactics as slide 3 (after first screenshot)
                if slide_index == 1 and competitor_tactics and len(competitor_tactics) > 0:
                    update_progress(progress + 5, "Adding competitor tactics...")
                    logger.info("Adding Competitor Tactics slide at position 3...")
                    self._add_competitor_tactics_slide(prs)

            # 4. Add any additional programmatic slides if needed
            update_progress(85, "Finalizing presentation...")

            # Save to bytes
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            file_size = len(output.getvalue())
            slide_count = len(prs.slides)
            logger.info("=" * 70)
            logger.info("POWERPOINT EXPORT WITH SCREENSHOTS COMPLETED")
            logger.info(f"Total slides: {slide_count}")
            logger.info(f"File size: {file_size:,} bytes ({file_size/1024:.1f} KB)")
            logger.info("=" * 70)

            return output.getvalue()

        except Exception as e:
            logger.error(f"Error in screenshot export: {e}")
            logger.info("Falling back to programmatic export...")
            return self.export_presentation(brand_name, industry, progress_callback)


def export_to_pptx(
    session_state: Dict[str, Any],
    brand_name: str = "Brand",
    industry: str = "General",
    progress_callback: Optional[callable] = None
) -> bytes:
    """
    Convenience function to export presentation (programmatic generation).
    """
    orchestrator = ExportOrchestrator(session_state)
    return orchestrator.export_presentation(
        brand_name=brand_name,
        industry=industry,
        progress_callback=progress_callback
    )


def export_to_pptx_with_screenshots(
    session_state: Dict[str, Any],
    brand_name: str = "Brand",
    industry: str = "General",
    app_url: str = "http://localhost:3006",
    use_live_capture: bool = True,
    progress_callback: Optional[callable] = None
) -> bytes:
    """
    Export presentation using actual screenshots of the Streamlit UI.

    Args:
        session_state: Streamlit session state dictionary
        brand_name: Brand name for the cover slide
        industry: Industry for the cover slide
        app_url: URL of the running Streamlit app (for live capture)
        use_live_capture: If True, capture from running app (pixel-perfect); if False, render HTML templates
        progress_callback: Optional callback for progress updates (percent, message)

    Returns:
        PowerPoint file bytes
    """
    orchestrator = ExportOrchestrator(session_state)
    return orchestrator.export_presentation_with_screenshots(
        brand_name=brand_name,
        industry=industry,
        app_url=app_url,
        use_live_capture=use_live_capture,
        progress_callback=progress_callback
    )
